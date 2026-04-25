"""Gradio entry point: Ask Finance — English UI, role-based Q&A, exports, logging."""

from __future__ import annotations

import ast
import json
import logging
import re
import sys
import uuid
from pathlib import Path
from typing import Any

import gradio as gr
import matplotlib
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.util import Pt

SRC_PATH = Path(__file__).resolve().parent / "src"
if SRC_PATH.exists() and str(SRC_PATH) not in sys.path:
    sys.path.insert(0, str(SRC_PATH))

from ask_finance import config
from ask_finance.agent import run_ask
from ask_finance.data_loaders import load_all, FinancialData
from ask_finance.logging_setup import setup_logging
from ask_finance.tools import (
    get_ebit_margin_trend,
    get_opex_variance,
    get_project_roi_trend,
)

logger = logging.getLogger("ask_finance.app")

matplotlib.use("Agg")

_fd: FinancialData | None = None
_LAST: dict[str, Any] = {}


def get_fd() -> FinancialData:
    global _fd
    if _fd is None:
        _fd = load_all()
    return _fd


def _roles() -> list[str]:
    return list(get_fd().rbac.get("roles", {}).keys()) or ["Group CFO"]


def _content_to_text(content: Any) -> str:
    """Normalize Gradio message content payloads into plain text."""
    if content is None:
        return ""
    if isinstance(content, str):
        txt = content.strip()
        # Recover from previously stringified message payloads like:
        # "[{'text': '...','type':'text'}]"
        for _ in range(4):
            if txt.startswith("[{") and "text" in txt:
                try:
                    parsed = ast.literal_eval(txt)
                except Exception:
                    break
                txt = _content_to_text(parsed).strip()
            else:
                break
        return txt
    if isinstance(content, dict):
        if "text" in content:
            return _content_to_text(content.get("text"))
        if "content" in content:
            return _content_to_text(content.get("content"))
        return json.dumps(content, ensure_ascii=False)
    if isinstance(content, list):
        chunks = [_content_to_text(x) for x in content]
        return "\n".join([x for x in chunks if x.strip()])
    return str(content)


def _normalize_history(history: list | None) -> list[dict[str, str]]:
    """
    Gradio Chatbot(type='messages') expects:
    [{"role": "user"|"assistant", "content": "..."}]
    Keep backward compatibility if tuple-style history appears.
    """
    if not history:
        return []
    normalized: list[dict[str, str]] = []
    for item in history:
        if isinstance(item, dict) and "role" in item and "content" in item:
            normalized.append(
                {
                    "role": str(item["role"]),
                    "content": _content_to_text(item["content"]),
                }
            )
            continue
        if (
            isinstance(item, (list, tuple))
            and len(item) == 2
        ):
            user_msg, assistant_msg = item
            normalized.append(
                {"role": "user", "content": _content_to_text(user_msg)}
            )
            normalized.append(
                {"role": "assistant", "content": _content_to_text(assistant_msg)}
            )
    return normalized


def _fallback_tool_trace(message: str, fd: FinancialData, role: str) -> list[dict[str, Any]]:
    """If model returns no function calls, run one deterministic tool by intent."""
    q = (message or "").lower()
    proj_years = sorted(fd.projects["reporting_year"].dropna().astype(int).unique().tolist())
    pl_years = sorted(fd.pl_monthly["fiscal_year"].dropna().astype(int).unique().tolist())
    latest_year = proj_years[-1] if proj_years else 2024
    start_default = latest_year - 2
    m_year = re.findall(r"\b(20\d{2})\b", q)
    if len(m_year) >= 2:
        start_y, end_y = int(m_year[0]), int(m_year[-1])
    elif len(m_year) == 1:
        end_y = int(m_year[0])
        start_y = end_y - 2
    else:
        start_y, end_y = start_default, latest_year

    if "roi" in q or "orion" in q:
        args = {"project_name": "Orion", "start_year": start_y, "end_year": end_y}
        return [
            {
                "tool": "get_project_roi_trend",
                "args": json.dumps(args),
                "result": get_project_roi_trend(fd, role, args),
                "s": 0.0,
            }
        ]

    q_match = re.search(r"\bq([1-4])\b", q)
    quarter = int(q_match.group(1)) if q_match else 2
    year_for_var = int(m_year[0]) if m_year else latest_year
    if "opex" in q and "variance" in q:
        bu = "Electronics" if "electronics" in q else None
        args = {"fiscal_year": year_for_var, "quarter": quarter}
        if bu:
            args["bu"] = bu
        return [
            {
                "tool": "get_opex_variance",
                "args": json.dumps(args),
                "result": get_opex_variance(fd, role, args),
                "s": 0.0,
            }
        ]

    if "ebit" in q and "margin" in q:
        if pl_years:
            if "all years" in q or "all year" in q:
                start_y = pl_years[0]
                end_y = pl_years[-1]
            else:
                start_y = max(start_y, pl_years[0])
                end_y = min(end_y, pl_years[-1])
                if start_y > end_y:
                    start_y, end_y = pl_years[0], pl_years[-1]
        args = {"start_year": start_y, "end_year": end_y}
        if "electronics" in q:
            args["bu"] = "Electronics"
        return [
            {
                "tool": "get_ebit_margin_trend",
                "args": json.dumps(args),
                "result": get_ebit_margin_trend(fd, role, args),
                "s": 0.0,
            }
        ]
    return []


def _build_insights_from_trace(trace: list[dict[str, Any]]) -> list[str]:
    insights: list[str] = []
    for step in trace:
        r = step.get("result") or {}
        if "by_fiscal_year" in r and r["by_fiscal_year"]:
            rows = r["by_fiscal_year"]
            if len(rows) >= 2:
                first = rows[0]
                last = rows[-1]
                delta = (last.get("ebit_margin", 0) - first.get("ebit_margin", 0)) * 100
                insights.append(
                    f"EBIT margin moved from {first.get('ebit_margin',0)*100:.2f}% ({first.get('fiscal_year')}) "
                    f"to {last.get('ebit_margin',0)*100:.2f}% ({last.get('fiscal_year')}), change {delta:+.2f} pts."
                )
            best = max(rows, key=lambda x: x.get("ebit_margin", 0))
            worst = min(rows, key=lambda x: x.get("ebit_margin", 0))
            insights.append(
                f"Best EBIT margin year: {best.get('fiscal_year')} ({best.get('ebit_margin',0)*100:.2f}%)."
            )
            insights.append(
                f"Weakest EBIT margin year: {worst.get('fiscal_year')} ({worst.get('ebit_margin',0)*100:.2f}%)."
            )
            break
        if "rows" in r and r["rows"] and "cumulative_roi" in r["rows"][0]:
            rows = r["rows"]
            first = rows[0]
            last = rows[-1]
            insights.append(
                f"Cumulative ROI changed from {first.get('cumulative_roi',0)*100:.2f}% ({first.get('reporting_year')}) "
                f"to {last.get('cumulative_roi',0)*100:.2f}% ({last.get('reporting_year')})."
            )
            break
    return insights


def answer_fn(message: str, history: list, role: str) -> tuple:
    safe_history = _normalize_history(history)
    if not (message or "").strip():
        return (
            safe_history,
            gr.update(),
            _LAST.get("trace_md", ""),
            _LAST.get("export_note", ""),
            _LAST.get("fig", None),
            gr.update(),
            gr.update(),
            _LAST.get("artifact_note", ""),
        )
    fd = get_fd()
    out = run_ask(fd, role, message.strip())
    text = out.get("answer", "")
    trace = out.get("tool_trace", [])
    if not trace:
        fallback = _fallback_tool_trace(message, fd, role)
        if fallback:
            trace = fallback
            text = (
                f"{text}\n\n(Used fallback finance tool execution for artifact generation.)"
            ).strip()
    trace_md = json.dumps(trace, indent=2) if trace else ""
    insight_lines = _build_insights_from_trace(trace)
    if "insight" in (message or "").lower() and insight_lines:
        if "Insight:" not in text and "Insights:" not in text:
            text = (
                f"{text}\n\nInsights:\n- "
                + "\n- ".join(insight_lines)
            ).strip()
    _LAST.clear()
    _LAST.update(
        {
            "tool_trace": trace,
            "answer": text,
            "role": role,
            "request_id": out.get("request_id"),
            "trace_md": trace_md,
            "insights": insight_lines,
        }
    )
    new_hist = safe_history + [
        {"role": "user", "content": _content_to_text(message)},
        {"role": "assistant", "content": _content_to_text(text)},
    ]
    export_note = f"Request {out.get('request_id')}; latency {out.get('latency_s')}s"
    # Simple chart: if last tool has by_fiscal_year, plot margin trend
    fig = None
    for step in reversed(trace):
        r = step.get("result") or {}
        if "by_fiscal_year" in r and r["by_fiscal_year"]:
            years = [x["fiscal_year"] for x in r["by_fiscal_year"]]
            margins = [x.get("ebit_margin", 0) for x in r["by_fiscal_year"]]
            fig, ax = plt.subplots(figsize=(6, 3))
            ax.plot(years, margins, marker="o")
            ax.set_xlabel("Fiscal year")
            ax.set_ylabel("EBIT margin")
            ax.set_title("EBIT margin trend (from tool output)")
            ax.grid(True, alpha=0.3)
            fig.tight_layout()
            break
        if "rows" in r and r["rows"]:
            x_key = "reporting_year"
            y_key = "cumulative_roi"
            if x_key in r["rows"][0] and y_key in r["rows"][0]:
                years = [x.get(x_key) for x in r["rows"]]
                roi_values = [x.get(y_key, 0) for x in r["rows"]]
                fig, ax = plt.subplots(figsize=(6, 3))
                ax.plot(years, roi_values, marker="o")
                ax.set_xlabel("Year")
                ax.set_ylabel("Cumulative ROI")
                ax.set_title("Project ROI trend (from tool output)")
                ax.grid(True, alpha=0.3)
                fig.tight_layout()
                break

    intent = (message or "").lower()
    wants_excel = any(k in intent for k in ("excel", "xlsx", "table", "sheet")) or bool(trace)
    wants_ppt = any(k in intent for k in ("ppt", "powerpoint", "slide", "deck", "summary")) or bool(trace)
    auto_excel = export_excel() if wants_excel else None
    auto_ppt = export_ppt() if wants_ppt else None
    artifact_note_parts: list[str] = []
    if fig is not None:
        artifact_note_parts.append("Chart generated")
    if auto_excel:
        artifact_note_parts.append("Excel ready")
    if auto_ppt:
        artifact_note_parts.append("PowerPoint ready")
    artifact_note = " | ".join(artifact_note_parts) if artifact_note_parts else "No artifact generated for this turn."
    _LAST["export_note"] = export_note
    _LAST["fig"] = fig
    _LAST["artifact_note"] = artifact_note

    return (
        new_hist,
        gr.update(value=""),
        trace_md,
        export_note,
        fig,
        gr.update(value=auto_excel) if auto_excel else gr.update(value=None),
        gr.update(value=auto_ppt) if auto_ppt else gr.update(value=None),
        artifact_note,
    )


def export_excel() -> str | None:
    trace = _LAST.get("tool_trace") or []
    Path(config.LOGS_DIR).mkdir(parents=True, exist_ok=True)
    path = config.LOGS_DIR / f"export_{_LAST.get('request_id', uuid.uuid4().hex[:8])}.xlsx"
    rows: list[dict] = []
    if trace:
        for i, step in enumerate(trace):
            r = step.get("result")
            if isinstance(r, dict):
                rows.append({"step": i, "tool": step.get("tool"), "summary": json.dumps(r)[:2000]})
            else:
                rows.append({"step": i, "tool": step.get("tool"), "summary": str(r)[:2000]})
    else:
        rows.append(
            {
                "step": 0,
                "tool": "none",
                "summary": (_LAST.get("answer") or "No tool trace for this turn.")[:2000],
            }
        )
    pd.DataFrame(rows).to_excel(path, index=False)
    return str(path)


def export_ppt() -> str | None:
    text = _LAST.get("answer")
    if not text:
        return None
    Path(config.LOGS_DIR).mkdir(parents=True, exist_ok=True)
    path = config.LOGS_DIR / f"summary_{_LAST.get('request_id', 'x')[:8]}.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Ask Finance — answer summary"
    tf = body.text_frame
    tf.text = (text or "")[:4000]
    for p in tf.paragraphs:
        p.font.size = Pt(12)
    insights = _LAST.get("insights") or []
    if insights:
        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Key insights"
        body2 = slide2.placeholders[1].text_frame
        body2.text = insights[0][:1000]
        for line in insights[1:5]:
            para = body2.add_paragraph()
            para.text = line[:1000]
            para.level = 0
        for p in body2.paragraphs:
            p.font.size = Pt(12)
    prs.save(path)
    return str(path)


def on_load():
    setup_logging()
    try:
        config.apply_credentials_env()
    except OSError as e:
        logger.warning("Credentials: %s", e)
    get_fd()
    logger.info("App load: data and logging ready. Project=%s", config.GOOGLE_PROJECT_ID)
    return None


with gr.Blocks(title="Ask Finance", theme=gr.themes.Soft()) as demo:
    gr.Markdown(
        """
# Ask Finance (demo)
English-only. Choose a **role** to simulate RBAC, then ask about P&L, opex variance, EBIT margin, or project ROI.
Example: *What was opex variance for Q2 2024 in the Electronics division?* — *Show ROI trend for Project Orion for 2021–2023.*
"""
    )
    role = gr.Dropdown(choices=_roles(), value="Group CFO", label="Role (RBAC simulation)")
    chat = gr.Chatbot(label="Chat", height=400)
    msg = gr.Textbox(label="Message", lines=2, placeholder="Your question in English…")
    plot = gr.Plot(label="Auto chart (when EBIT margin trend tool is used)")
    with gr.Accordion("Debug (collapsed by default)", open=False):
        trace = gr.Textbox(label="Tool trace (JSON)", lines=8)
        status = gr.Textbox(label="Last run", lines=1)

    with gr.Accordion("Artifacts (collapsed by default)", open=False):
        artifact_note = gr.Textbox(label="Artifact status", lines=1)
        with gr.Row():
            b_x = gr.Button("Export last run to Excel")
            b_p = gr.Button("Export last answer to PowerPoint")
        with gr.Row():
            f_x = gr.File(label="Download Excel")
            f_p = gr.File(label="Download PPT")

    def submit(m, h, r):
        return answer_fn(m, h, r)

    msg.submit(
        submit,
        [msg, chat, role],
        [chat, msg, trace, status, plot, f_x, f_p, artifact_note],
    )
    b_send = gr.Button("Send")
    b_send.click(
        submit,
        [msg, chat, role],
        [chat, msg, trace, status, plot, f_x, f_p, artifact_note],
    )

    def do_x():
        p = export_excel()
        return gr.update(value=p) if p else gr.update()

    def do_p():
        p = export_ppt()
        return gr.update(value=p) if p else gr.update()

    b_x.click(do_x, outputs=[f_x])
    b_p.click(do_p, outputs=[f_p])

    demo.load(on_load)

if __name__ == "__main__":
    setup_logging()
    config.apply_credentials_env()
    get_fd()
    demo.queue().launch(server_name="0.0.0.0", server_port=7860, share=False)
